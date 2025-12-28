import os
import platform
import subprocess
import tempfile
import shutil
import zipfile

#libreoffice도 다운로드 필요해요 ~.~

import pdfplumber
from docx import Document
from shutil import which

DATA_DIR = "data"  # 크롤러랑 맞추기
FILE_DIR = os.path.join(DATA_DIR, "files")
FILE_TEXT_DIR = os.path.join(DATA_DIR, "file_text")

os.makedirs(FILE_TEXT_DIR, exist_ok=True)

OS_TYPE = platform.system()  # "Windows" / "Linux" / "Darwin"

# 윈도우에서 자주 쓰는 기본 설치 경로 후보
SOFFICE_CANDIDATES = [
    "soffice",
    "libreoffice",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
]


# ========================
# 공통 유틸
# ========================
def has_libreoffice() -> bool:
    """시스템에 libreoffice/soffice가 있는지 확인 (PATH + 기본 설치 경로 둘 다 확인)"""
    # 1) PATH 에 있는지 확인
    if which("libreoffice") is not None or which("soffice") is not None:
        return True

    # 2) 윈도우 기본 설치 경로 확인
    for p in SOFFICE_CANDIDATES:
        if os.path.isabs(p) and os.path.exists(p):
            return True

    return False

from typing import Optional
def _find_soffice_exe() -> Optional[str]:
    """실행 가능한 soffice / libreoffice 경로 찾기"""
    # 1) PATH 검색
    for name in ["libreoffice", "soffice"]:
        exe = which(name)
        if exe is not None:
            return exe

    # 2) 우리가 아는 기본 경로 검색
    for p in SOFFICE_CANDIDATES:
        if os.path.isabs(p) and os.path.exists(p):
            return p

    return None


def run_libreoffice_convert(src_path: str, target_ext: str, out_dir: str) -> str:
    """
    LibreOffice로 파일 변환.
    예: doc -> docx, ppt -> pptx, xls -> xlsx 등.
    성공하면 새 파일 경로, 실패하면 원본 경로 반환.
    """
    exe = _find_soffice_exe()
    if exe is None:
        print(
            "[FAIL] LibreOffice 실행 파일을 찾지 못했습니다. PATH 또는 설치 경로를 확인하세요."
        )
        return src_path

    cmd = [
        exe,
        "--headless",
        "--convert-to",
        target_ext,
        src_path,
        "--outdir",
        out_dir,
    ]

    try:
        proc = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            check=False,
        )
        # 디버깅하고 싶으면 아래 두 줄 한 번 찍어봐도 됨
        # print("[LIBREOFFICE STDOUT]", proc.stdout)
        # print("[LIBREOFFICE STDERR]", proc.stderr)
    except Exception as e:
        print("[FAIL] LibreOffice 실행 실패:", e)
        return src_path

    base = os.path.splitext(os.path.basename(src_path))[0]
    out_path = os.path.join(out_dir, base + "." + target_ext)

    if os.path.exists(out_path):
        print(f"[CONVERT] {src_path} → {out_path}")
        return out_path
    else:
        print("[FAIL] LibreOffice 변환 실패:", src_path)
        return src_path


# ========================
# 확장자별 텍스트 추출 함수
# ========================
def extract_pdf(path: str) -> str:
    texts = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                texts.append(page.extract_text() or "")
    except Exception as e:
        print("[SKIP] PDF 추출 실패:", path, e)
        return ""
    return "\n".join(texts)


def extract_docx(path: str) -> str:
    try:
        doc = Document(path)
    except Exception as e:
        print("[SKIP] DOCX 열기 실패:", path, e)
        return ""

    texts = []

    # 1) 일반 문단 텍스트
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            texts.append(p.text.strip())

    # 2) 표(Table) 안 텍스트
    for table in doc.tables:
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    row_cells.append(t)
            if row_cells:
                # 한 줄로 합치거나, 탭으로 구분하고 싶으면 "\t".join으로 변경 가능
                texts.append(" | ".join(row_cells))

    if not texts:
        print("[INFO] DOCX에서 유의미한 텍스트를 찾지 못함:", path)

    return "\n".join(texts)


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(path)
    except Exception as e:
        print(f"[SKIP] PPTX 열기 실패: {path} -> {e}")
        return ""

    texts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # 1) 일반 텍스트
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if shape.text:
                    texts.append(shape.text)

            # 2) 표 텍스트
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(cell.text)

    return "\n".join(t.strip() for t in texts if t and t.strip())


def extract_xlsx(path: str) -> str:
    try:
        import openpyxl
    except ImportError:
        print("[SKIP] openpyxl 미설치:", path)
        return ""

    try:
        wb = openpyxl.load_workbook(path, data_only=True)
    except Exception as e:
        print("[SKIP] XLSX 열기 실패:", path, e)
        return ""

    texts = []
    for sheet in wb.worksheets:
        texts.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            # 각 셀을 문자열로 변환 후 탭으로 join
            row_str = "\t".join("" if v is None else str(v) for v in row)
            texts.append(row_str)

    return "\n".join(texts)


# ========================
# Windows 전용 HWP/HWPX (수정 버전)
# ========================
def extract_hwp_windows_only(path: str) -> str:
    if OS_TYPE != "Windows":
        print("[SKIP] HWP: Windows 아님:", path)
        return ""

    try:
        import win32com.client as win32
    except ImportError:
        print("[SKIP] HWP: pywin32 미설치:", path)
        return ""

    # FILE_TEXT_DIR 안에 저장할 txt 경로 (절대 경로로!)
    base = os.path.splitext(os.path.basename(path))[0]
    out_txt_rel = os.path.join(FILE_TEXT_DIR, base + ".txt")
    out_txt = os.path.abspath(out_txt_rel)  # 여기서 절대 경로로 변환

    # 혹시 폴더가 없다면 생성
    os.makedirs(os.path.dirname(out_txt), exist_ok=True)

    try:
        hwp = win32.gencache.EnsureDispatch("HWPFrame.HwpObject")

        try:
            hwp.RegisterModule("FilePathCheckDLL", "SecurityModule")
        except Exception:
            pass

        # 파일 열기
        hwp.Open(path, "", "forceopen:true")

        # 텍스트 파일로 저장 (절대 경로 사용)
        hwp.SaveAs(out_txt, "TEXT")

        # 문서 닫기
        hwp.XHwpDocuments.Close(isDirty=False)

        # 저장된 TXT 읽기
        try:
            with open(out_txt, "r", encoding="cp949", errors="ignore") as f:
                return f.read()
        except:
            with open(out_txt, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

    except Exception as e:
        print("[SKIP] HWP 변환 실패:", path, e)
        return ""
    finally:
        try:
            hwp.Quit()
        except Exception:
            pass


# ========================
# Windows 전용 DOC/PPT (COM)
# ========================
def extract_doc_windows_com(path: str) -> str:
    if OS_TYPE != "Windows":
        return ""

    try:
        import win32com.client as win32
    except ImportError:
        return ""

    try:
        word = win32.gencache.EnsureDispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(path)
        text = doc.Content.Text
        doc.Close(False)
        word.Quit()
        return text
    except Exception as e:
        print("[SKIP] DOC COM 추출 실패:", path, e)
        try:
            word.Quit()
        except Exception:
            pass
        return ""
    
# ========================
# Windows 전용 HWP/HWPX
# ========================
def extract_hwp_windows_only(path: str) -> str:
    if OS_TYPE != "Windows":
        print("[SKIP] HWP: Windows 아님:", path)
        return ""

    try:
        import win32com.client as win32
    except ImportError:
        print("[SKIP] HWP: pywin32 미설치:", path)
        return ""

    # 절대 경로로 변환
    src_path = os.path.abspath(path)

    base_name = os.path.splitext(os.path.basename(path))[0]
    out_txt = os.path.join(FILE_TEXT_DIR, base_name + ".txt")
    out_txt = os.path.abspath(out_txt)

    os.makedirs(os.path.dirname(out_txt), exist_ok=True)

    try:
        hwp = win32.gencache.EnsureDispatch("HWPFrame.HwpObject")

        # [중요] 일단 보안 모듈 등록은 빼고 테스트
        # try:
        #     hwp.RegisterModule("FilePathCheckDLL", "SecurityModule")
        # except Exception:
        #     pass

        # hwp_to_txt_batch.py와 똑같이 두 번째 인수는 ""
        hwp.Open(src_path, "", "forceopen:true")

        # 바로 텍스트로 저장
        hwp.SaveAs(out_txt, "TEXT")

        # 문서 닫기
        try:
            hwp.XHwpDocuments.Close(isDirty=False)
        except Exception:
            pass

        # 한글 종료
        try:
            hwp.Quit()
        except Exception:
            pass

    except Exception as e:
        print("[SKIP] HWP 변환 실패:", src_path, e)
        try:
            hwp.Quit()
        except Exception:
            pass
        return ""

    # 실제 파일 크기 확인
    if not os.path.exists(out_txt):
        print("[SKIP] HWP TXT 파일이 생성되지 않음:", out_txt)
        return ""

    size = os.path.getsize(out_txt)
    if size == 0:
        print("[SKIP] HWP TXT 파일 크기 0:", out_txt)
        return ""

    # 텍스트 읽어서 반환
    try:
        with open(out_txt, "r", encoding="cp949", errors="ignore") as f:
            text = f.read()
    except UnicodeDecodeError:
        with open(out_txt, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

    if not text.strip():
        print("[SKIP] HWP TXT 내용 비어 있음:", out_txt)
        return ""

    print("[HWP] 텍스트 추출 성공:", out_txt)
    return text


# ========================
# ZIP 처리
# ========================
def extract_zip(path: str) -> str:
    """
    zip 하나를 '컨테이너'로 보고,
    내부 파일마다 개별 txt 파일을 생성합니다.

    반환값은 "" (zip 자체에 대한 텍스트는 없음)
    """
    tmp_dir = tempfile.mkdtemp(prefix="zip_extract_")

    try:
        with zipfile.ZipFile(path, "r") as zf:
            members = [m for m in zf.namelist() if not m.endswith("/")]

            if not members:
                print("[SKIP] ZIP 안에 파일이 없음:", path)
                return ""

            zip_base = os.path.splitext(os.path.basename(path))[0]

            for name in members:
                extracted_path = zf.extract(name, path=tmp_dir)

                # 이름 정리: zip이름__원본파일이름.txt 이런 식으로 저장
                inner_base = os.path.splitext(os.path.basename(name))[0]
                out_txt_name = f"{zip_base}__{inner_base}.txt"
                out_txt_path = os.path.join(FILE_TEXT_DIR, out_txt_name)

                # 이미 만들어진 적 있으면 스킵
                if os.path.exists(out_txt_path):
                    print("[SKIP] 이미 존재:", out_txt_path)
                    continue

                # 내부 파일에 대해 다시 extract_by_ext 호출
                inner_text = extract_by_ext(extracted_path)

                if not inner_text.strip():
                    print(f"[SKIP] ZIP 내부 파일 텍스트 없음: {name}")
                    continue

                with open(out_txt_path, "w", encoding="utf-8", errors="ignore") as f:
                    f.write(inner_text)

                print(f"[ZIP] {name} → {out_txt_path}")

    except Exception as e:
        print("[SKIP] ZIP 처리 실패:", path, e)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    # zip 자체로는 텍스트를 리턴하지 않음
    return ""


# ========================
# 확장자 라우팅
# ========================
def extract_by_ext(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()

    # PDF
    if ext == ".pdf":
        return extract_pdf(path)

    # 최신 office
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".xlsx":
        return extract_xlsx(path)

    # 구형 office (.doc / .ppt / .xls) → OS 상관 없이 LibreOffice로만 처리
    if ext == ".doc":
        if has_libreoffice():
            tmp_dir = tempfile.mkdtemp(prefix="doc_convert_")
            try:
                new_file = run_libreoffice_convert(path, "docx", tmp_dir)
                if new_file != path and os.path.exists(new_file):
                    return extract_docx(new_file)
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)
        print("[SKIP] DOC 처리 불가 (LibreOffice 없음):", path)
        return ""

    if ext == ".ppt":
        if has_libreoffice():
            tmp_dir = tempfile.mkdtemp(prefix="ppt_convert_")
            try:
                new_file = run_libreoffice_convert(path, "pptx", tmp_dir)
                if new_file != path and os.path.exists(new_file):
                    return extract_pptx(new_file)
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)
        print("[SKIP] PPT 처리 불가 (LibreOffice 없음):", path)
        return ""

    if ext == ".xls":
        if has_libreoffice():
            tmp_dir = tempfile.mkdtemp(prefix="xls_convert_")
            try:
                new_file = run_libreoffice_convert(path, "xlsx", tmp_dir)
                if new_file != path and os.path.exists(new_file):
                    return extract_xlsx(new_file)
            finally:
                shutil.rmtree(tmp_dir, ignore_errors=True)
        print("[SKIP] XLS 처리 불가 (LibreOffice 없음):", path)
        return ""

    # HWP/HWPX
    if ext in [".hwp", ".hwpx"]:
        return extract_hwp_windows_only(path)

    # ZIP
    if ext == ".zip":
        return extract_zip(path)

    print("[SKIP] 지원하지 않는 확장자:", path)
    return ""


# ========================
# 메인(단독 실행 시)
# ========================
def main():
    file_count = 0
    text_count = 0

    if not os.path.isdir(FILE_DIR):
        print("[WARN] FILE_DIR가 존재하지 않습니다:", FILE_DIR)
        return

    for fname in os.listdir(FILE_DIR):
        src = os.path.join(FILE_DIR, fname)
        if not os.path.isfile(src):
            continue

        base = os.path.splitext(fname)[0]
        out_txt = os.path.join(FILE_TEXT_DIR, base + ".txt")

        if os.path.exists(out_txt):
            continue

        file_count += 1
        print(f"\n[{file_count}] 처리중: {src}")

        text = extract_by_ext(src)
        if not text.strip():
            print(" → 텍스트 없음")
            continue

        with open(out_txt, "w", encoding="utf-8", errors="ignore") as f:
            f.write(text)

        text_count += 1
        print(" → 텍스트 저장:", out_txt)

    print("\n[DONE] 파일 수:", file_count)
    print("[DONE] 텍스트 추출된 파일 수:", text_count)


if __name__ == "__main__":
    main()
